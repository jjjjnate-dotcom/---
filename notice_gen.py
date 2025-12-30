#!/usr/bin/env python
"""
안내문 작성 비서 (notice_gen.py) - A4 안내문 PPT 생성기
필요: pip install python-pptx
실행 예시:
  - 템플릿 만들기: python notice_gen.py --template -o notice_template.pptx
  - 텍스트 파일 본문 적용: python notice_gen.py --body-text body.txt -o notice_a4.pptx
  - JSON 데이터 적용: python notice_gen.py --data data.json -o notice_a4.pptx
  - 터미널 입력 페이지: python notice_gen.py --interactive -o notice_a4.pptx --export-json meta.json --open
  - GUI 입력 페이지: python notice_gen.py --gui
  - 웹 입력 페이지 서버: python notice_gen.py --serve --port 8000 (브라우저로 notice_form.html 또는 http://localhost:8000 접속)
"""
import argparse
import json
from pathlib import Path
import os
import tkinter as tk
from tkinter import filedialog, messagebox
from io import BytesIO
from http.server import BaseHTTPRequestHandler, HTTPServer

from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

A4_W, A4_H = Mm(210), Mm(297)
PRINT_MARGIN = Mm(12)  # 인쇄 안전 여백
FONT_NAME = "Malgun Gothic"
BLUE = RGBColor(30, 115, 200)   # 메인 색
DARK = RGBColor(12, 74, 153)    # 포인트
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(29, 29, 31)
LIGHT = RGBColor(236, 242, 249)
OUTLINE = RGBColor(205, 217, 231)


def parse_period(raw):
    parts = raw.split("~")
    if len(parts) == 2:
        start_v = parts[0].strip().replace(".", "-")
        end_v = parts[1].strip().replace(".", "-")
        return start_v or "YYYY-MM-DD", end_v or "YYYY-MM-DD"
    return ("YYYY-MM-DD", "YYYY-MM-DD")


def build_data(notice_no, apt, period, title, body_lines=None):
    start_v, end_v = parse_period(period) if isinstance(period, str) else ("YYYY-MM-DD", "YYYY-MM-DD")
    return {
        "title": title or "제목을 입력하세요",
        "label": "게시기간",
        "start": start_v,
        "end": end_v,
        "notice_no": notice_no,
        "body": body_lines if body_lines is not None else ["(AI 본문 자리)"],
        "footer": apt or "발신처를 입력하세요 (예: 000아파트 관리사무소장 [직인생략])",
    }


def add_box(slide, left, top, width, height, fill, text="", size=18, bold=True, color=WHITE, align=PP_ALIGN.LEFT, shadow=False):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if shadow:
        shape.shadow.inherit = False
        shape.shadow.visible = True
        shape.shadow.blur_radius = Pt(6)
        shape.shadow.distance = Pt(1.5)
        shape.shadow.direction = 270
    if text:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.name = FONT_NAME
        p.font.color.rgb = color
        p.alignment = align
    return shape


def normalize_body(body):
    if body is None:
        return []
    if isinstance(body, str):
        return [line.rstrip() for line in body.splitlines()] or [body]
    return [str(line) for line in body]


def load_body_from_text(path):
    text = Path(path).read_text(encoding="utf-8")
    lines = [line.rstrip() for line in text.splitlines()]
    return lines or [""]


def make_notice(data, filename="notice_a4.pptx"):
    prs = Presentation()
    prs.slide_width, prs.slide_height = A4_W, A4_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    header_h = Mm(18)  # 하단 푸터와 동일한 높이
    inner_w = A4_W - PRINT_MARGIN * 2

    header = add_box(slide, PRINT_MARGIN, PRINT_MARGIN, inner_w, header_h, BLUE,
            text=data["title"], size=24, align=PP_ALIGN.CENTER, shadow=True)

    info_bar_h = Mm(12)
    info_top = PRINT_MARGIN + header_h  # 제목 칸에 바로 밀착
    info_bar = slide.shapes.add_shape(1, PRINT_MARGIN, info_top, inner_w, info_bar_h)
    info_bar.fill.solid()
    info_bar.fill.fore_color.rgb = DARK
    info_bar.line.fill.background()
    tf_info = info_bar.text_frame
    tf_info.clear()
    p_info = tf_info.paragraphs[0]
    notice_no_txt = f"공고번호: {data['notice_no']}" if data.get("notice_no") else "공고번호: -"
    date_txt = f"{data['start']} ~ {data['end']}"
    p_info.text = f"{notice_no_txt}   |   게시기간: {date_txt}"
    p_info.font.size = Pt(12)
    p_info.font.name = FONT_NAME
    p_info.font.color.rgb = WHITE
    p_info.alignment = PP_ALIGN.CENTER

    body_top = info_top + info_bar_h + Mm(6)  # 위로 당겨진 만큼 본문 영역도 커짐
    body_h = A4_H - PRINT_MARGIN - body_top - Mm(24)

    # 본문 영역 배경 박스 (연한 톤, 외곽선)
    body_bg = slide.shapes.add_shape(1, PRINT_MARGIN, body_top, inner_w, body_h)
    body_bg.fill.solid()
    body_bg.fill.fore_color.rgb = LIGHT
    body_bg.line.color.rgb = OUTLINE
    body_bg.line.width = Pt(0.75)
    body = slide.shapes.add_textbox(PRINT_MARGIN + Mm(4), body_top + Mm(4), inner_w - Mm(8), body_h - Mm(8))
    tf = body.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Mm(2)
    tf.word_wrap = True

    for i, line in enumerate(normalize_body(data["body"])):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(14)
        p.font.name = FONT_NAME
        p.font.color.rgb = TEXT
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = 1.0

    footer_h = Mm(18)
    add_box(slide, PRINT_MARGIN, A4_H - PRINT_MARGIN - footer_h, inner_w, footer_h, BLUE,
            text=data["footer"], size=15, align=PP_ALIGN.CENTER, bold=True, shadow=True)

    prs.save(filename)
    print(f"saved: {filename}")


def make_template(filename="notice_template.pptx", body_lines=None):
    template_data = {
        "title": "제목을 입력하세요",
        "label": "게시기간",
        "start": "YYYY-MM-DD",
        "end": "YYYY-MM-DD",
        "notice_no": "공고번호를 입력하세요",
        "body": body_lines if body_lines is not None else ["여기에 본문을 입력하세요."],
        "footer": "발신처를 입력하세요 (예: 관리사무소장 [직인생략])",
    }
    make_notice(template_data, filename)


def parse_args():
    parser = argparse.ArgumentParser(description="A4 안내문 PPT 생성기")
    parser.add_argument("--data", help="title/label/dates/body/footer가 담긴 JSON 파일 경로")
    parser.add_argument("--body-text", help="본문을 줄바꿈 기준으로 적은 UTF-8 텍스트 파일 경로")
    parser.add_argument("--interactive", action="store_true", help="터미널에서 제목/게시기간/공고번호/푸터를 입력")
    parser.add_argument("--gui", action="store_true", help="간단한 GUI 입력 폼으로 제목/공고번호/게시기간/단지명 입력")
    parser.add_argument("--serve", action="store_true", help="로컬 웹 서버 실행 후 notice_form.html로 입력 받아 PPTX 바로 다운로드")
    parser.add_argument("--port", type=int, default=8000, help="--serve 모드 포트 (기본 8000)")
    parser.add_argument("--export-json", help="입력 받은 메타데이터를 JSON으로 저장할 경로")
    parser.add_argument("--template", action="store_true", help="입력용 템플릿 PPTX를 생성")
    parser.add_argument("--open", action="store_true", help="생성 후 PPTX 파일을 바로 열기")
    parser.add_argument("-o", "--output", default=None, help="저장할 pptx 파일명")
    return parser.parse_args()


def prompt_meta():
    print("안내문 메타데이터를 입력하세요. (엔터 시 기본값 적용)")
    title = input("제목: ").strip() or "제목을 입력하세요"
    start = input("게시 시작일(YYYY-MM-DD): ").strip() or "YYYY-MM-DD"
    end = input("게시 종료일(YYYY-MM-DD): ").strip() or "YYYY-MM-DD"
    notice_no = input("공고번호: ").strip() or ""
    footer = input("푸터/발신처: ").strip() or "발신처를 입력하세요 (예: 관리사무소장 [직인생략])"
    return {
        "title": title,
        "label": "게시기간",
        "start": start,
        "end": end,
        "notice_no": notice_no,
        "body": ["(AI 본문 자리)"],
        "footer": footer,
    }


def gui_collect():
    root = tk.Tk()
    root.title("안내문 메타데이터 입력")
    root.geometry("480x320")
    root.resizable(False, False)

    labels = ["공고번호", "아파트 이름", "게시기간 (예: 2025.12.29 ~ 2026.01.05)", "제목"]
    defaults = [
        "종합S(갑) 제2512-001호",
        "한공원 아파트",
        "2025.12.29 ~ 2026.01.05",
        "층간소음 안내문",
    ]
    entries = []

    for i, (label, default) in enumerate(zip(labels, defaults)):
        tk.Label(root, text=label, anchor="w").grid(row=i, column=0, sticky="w", padx=12, pady=6)
        entry = tk.Entry(root, width=50)
        entry.insert(0, default)
        entry.grid(row=i, column=1, padx=12, pady=6)
        entries.append(entry)

    open_var = tk.BooleanVar(value=True)
    result = {}

    def on_submit():
        notice_v, apt_v, period_v, title_v = [e.get().strip() for e in entries]
        start_v, end_v = parse_period(period_v)
        if not title_v:
            messagebox.showerror("입력 오류", "제목은 필수입니다.")
            return
        out_path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint", "*.pptx")],
            initialfile="notice_a4.pptx",
            title="안내문 저장 위치 선택",
        )
        if not out_path:
            return
        result["data"] = {
            "title": title_v,
            "label": "게시기간",
            "start": start_v or "YYYY-MM-DD",
            "end": end_v or "YYYY-MM-DD",
            "notice_no": notice_v,
            "body": ["(AI 본문 자리)"],
            "footer": apt_v or "발신처를 입력하세요 (예: 000아파트 관리사무소장 [직인생략])",
        }
        result["output"] = out_path
        result["open"] = open_var.get()
        root.destroy()

    tk.Checkbutton(root, text="생성 후 바로 열기", variable=open_var).grid(row=len(labels), column=0, columnspan=2, pady=4)
    btn = tk.Button(root, text="PPTX 생성", command=on_submit, width=20)
    btn.grid(row=len(labels)+1, column=0, columnspan=2, pady=10)

    root.mainloop()
    return result.get("data"), result.get("output"), result.get("open")


def open_file(path):
    try:
        os.startfile(path)
    except Exception as e:
        print(f"파일을 자동으로 열지 못했습니다: {e}")


def run_server(port):
    class Handler(BaseHTTPRequestHandler):
        def _send(self, status=200, content_type="text/plain", body=b""):
            self.send_response(status)
            self.send_header("Content-Type", content_type)
            self.send_header("Access-Control-Allow-Origin", "*")
            self.end_headers()
            if body:
                self.wfile.write(body)

        def do_OPTIONS(self):
            self.send_response(200)
            self.send_header("Access-Control-Allow-Origin", "*")
            self.send_header("Access-Control-Allow-Methods", "GET,POST,OPTIONS")
            self.send_header("Access-Control-Allow-Headers", "Content-Type")
            self.end_headers()

        def do_GET(self):
            if self.path == "/":
                html_path = Path(__file__).with_name("notice_form.html")
                if not html_path.exists():
                    self._send(404, "text/plain", b"notice_form.html not found")
                    return
                content = html_path.read_bytes()
                self._send(200, "text/html; charset=utf-8", content)
            else:
                self._send(404, "text/plain", b"Not found")

        def do_POST(self):
            if self.path != "/generate":
                self._send(404, "text/plain", b"Not found")
                return
            length = int(self.headers.get("Content-Length", 0))
            raw = self.rfile.read(length) if length else b"{}"
            try:
                payload = json.loads(raw.decode("utf-8"))
            except Exception:
                self._send(400, "text/plain", b"Invalid JSON")
                return

            notice_no = payload.get("notice_no", "")
            apt_name = payload.get("apt_name", "")
            period = payload.get("period", "YYYY-MM-DD ~ YYYY-MM-DD")
            title = payload.get("title", "제목을 입력하세요")
            body_lines = payload.get("body")
            data = build_data(notice_no, apt_name, period, title, body_lines)

            stream = BytesIO()
            make_notice(data, filename=stream)
            pptx_bytes = stream.getvalue()

            self.send_response(200)
            self.send_header("Content-Type", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
            self.send_header("Content-Disposition", 'attachment; filename="notice_a4.pptx"')
            self.send_header("Access-Control-Allow-Origin", "*")
            self.send_header("Content-Length", str(len(pptx_bytes)))
            self.end_headers()
            self.wfile.write(pptx_bytes)

    server = HTTPServer(("0.0.0.0", port), Handler)
    print(f"안내문 작성 비서 서버 실행 중: http://localhost:{port}")
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        pass
    finally:
        server.server_close()


if __name__ == "__main__":
    args = parse_args()

    if args.serve:
        run_server(args.port)
        raise SystemExit

    default_output = "notice_template.pptx" if args.template else "notice_a4.pptx"
    output = args.output or default_output

    if args.template:
        body_lines = load_body_from_text(args.body_text) if args.body_text else None
        make_template(filename=output, body_lines=body_lines)
    else:
        if args.gui:
            data, gui_output, gui_open = gui_collect()
            if not data:
                raise SystemExit("GUI 입력이 취소되었습니다.")
            if gui_output:
                output = gui_output
            if gui_open:
                args.open = True
        elif args.interactive:
            data = prompt_meta()
        else:
            data = {
                "title": "반려견 목줄 착용안내",
                "label": "게시기간",
                "start": "2025-10-14",
                "end": "2025-10-21",
                "notice_no": "제2025-001호",
                "body": [
                    "관리사무실에서 안내드립니다.",
                    "공동주택 내에서는 반려견 목줄 착용이 의무사항입니다.",
                    "다른 입주민의 안전과 불쾌감 방지를 위해 주의해주시기 바랍니다.",
                    "1. 반려견은 외출 시 반드시 목줄을 착용해 주세요.",
                    "2. 엘리베이터나 계단 이용 시에도 꼭 붙잡아 주시기 바랍니다.",
                    "3. 어린이 놀이터, 공용공간에서의 출입은 자제해 주시기 바랍니다.",
                    "4. 배설물은 즉시 수거해 주시고, 위생에도 신경 써 주세요.",
                    "5. 지속적인 민원이 발생하면 관련 규정에 따라 조치가 있을 수 있습니다.",
                    "사소해 보여도 이웃 간 불편함이 생길 수 있는 부분입니다.",
                    "함께 사는 공간인 만큼, 기본적인 예의를 지켜주세요.",
                    "협조해 주셔서 감사합니다."
                ],
                "footer": "인계수정아파트관리사무소장 관리사무소장 [직인생략]"
            }

        if args.data:
            data = json.loads(Path(args.data).read_text(encoding="utf-8"))
        if args.body_text:
            data["body"] = load_body_from_text(args.body_text)

        if args.export_json:
            Path(args.export_json).write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")

        make_notice(data, filename=output)
        if args.open:
            open_file(output)
